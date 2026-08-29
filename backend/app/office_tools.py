"""办公杂活：文字→PPT大纲、图片→PPT、工资条拆分、多表合并去重。
全部纯本地（python-pptx / openpyxl / Pillow），不调任何模型。

PPT→PDF 不在这里——那个要 LibreOffice 渲染引擎，是另一个服务器依赖决定。
"""
import io
import zipfile

from openpyxl import Workbook, load_workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

from app import doc_format


# ---------------------------------------------------------------------------
# 文字 / 大纲 -> PPT
# ---------------------------------------------------------------------------

def _set_cjk(run, name: str = "微软雅黑") -> None:
    """python-pptx 的 run.font.name 只设 <a:latin>，中文要单独设 <a:ea>。"""
    from pptx.oxml.ns import qn

    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", name)


def text_to_pptx(text: str, title: str | None = None) -> bytes:
    """把结构化文字转成 PPT 大纲：一级标题=一页、其下的段落/条目=该页要点。
    没有标题层级时，首行做封面标题，其余按段落分页（每页≤6条）。"""
    blocks = doc_format.parse(text, explicit_title=(title or None))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    doc_title = ""
    heads = [b for b in blocks if b.kind == "heading"]

    # 封面
    tb = next((b for b in blocks if b.kind == "title"), None)
    if tb:
        doc_title = doc_format._plain(tb)
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = doc_title
        for r in s.shapes.title.text_frame.paragraphs[0].runs:
            _set_cjk(r)
        if len(s.placeholders) > 1:
            s.placeholders[1].text = ""

    def new_content_slide(heading_text: str):
        s = prs.slides.add_slide(content_layout)
        s.shapes.title.text = heading_text
        for r in s.shapes.title.text_frame.paragraphs[0].runs:
            _set_cjk(r)
        body = s.placeholders[1].text_frame
        body.clear()
        return body, s

    if heads:
        cur_body = None
        first = True
        for b in blocks:
            if b.kind in ("title",):
                continue
            if b.kind == "heading" and b.level <= 2:
                cur_body, _ = new_content_slide(doc_format._plain(b))
                first = True
                continue
            txt = doc_format._plain(b).strip()
            if not txt or cur_body is None:
                continue
            para = cur_body.paragraphs[0] if first else cur_body.add_paragraph()
            first = False
            para.text = txt
            para.level = 1 if b.kind in ("heading",) else 0
            for r in para.runs:
                _set_cjk(r)
                r.font.size = Pt(20)
    else:
        # 无标题层级：按段落切片
        paras = [doc_format._plain(b).strip() for b in blocks
                 if b.kind in ("para", "bullet", "ordered") and doc_format._plain(b).strip()]
        per = 6
        for i in range(0, len(paras), per):
            body, _ = new_content_slide(doc_title or "内容" if i == 0 else "（续）")
            for k, line in enumerate(paras[i:i + per]):
                para = body.paragraphs[0] if k == 0 else body.add_paragraph()
                para.text = line
                for r in para.runs:
                    _set_cjk(r)
                    r.font.size = Pt(20)

    if len(prs.slides) == 0:
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = doc_title or "（空）"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def images_to_pptx(images: list[bytes]) -> bytes:
    """每张图片一页，等比缩放居中铺在 16:9 页面上。"""
    from PIL import Image as PILImage, ImageOps

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    pw, ph = prs.slide_width, prs.slide_height

    for raw in images:
        im = PILImage.open(io.BytesIO(raw))
        im = ImageOps.exif_transpose(im).convert("RGB")
        if max(im.size) > 2200:
            s = 2200 / max(im.size)
            im = im.resize((round(im.width * s), round(im.height * s)), PILImage.LANCZOS)
        b = io.BytesIO()
        im.save(b, "JPEG", quality=88)
        b.seek(0)

        scale = min(pw / im.width, ph / im.height)
        w, h = int(im.width * scale), int(im.height * scale)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(b, Emu(int((pw - w) / 2)), Emu(int((ph - h) / 2)), width=Emu(w), height=Emu(h))

    if len(prs.slides) == 0:
        prs.slides.add_slide(blank)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# 工资条拆分
# ---------------------------------------------------------------------------

_THIN = Side(style="thin", color="808080")
_BORDER = Border(left=_THIN, right=_THIN, top=_THIN, bottom=_THIN)
_HEAD_FILL = PatternFill("solid", fgColor="F2F2F2")


def _read_first_sheet(data: bytes):
    wb = load_workbook(io.BytesIO(data), data_only=True)
    ws = wb[wb.sheetnames[0]]
    rows = [[("" if c is None else c) for c in row] for row in ws.iter_rows(values_only=True)]
    # 去掉完全空的尾行
    while rows and all(str(c).strip() == "" for c in rows[-1]):
        rows.pop()
    return rows


def payslips(data: bytes, slip_title: str = "", per_page: int = 12) -> bytes:
    """一张工资总表 -> 每人一条工资条，堆在一个 sheet 里，设好打印区域，打印后裁开即可。
    每条 = 表头行 + 该员工数据行（+ 可选标题行）。"""
    rows = _read_first_sheet(data)
    if len(rows) < 2:
        raise ValueError("表格至少要有表头行 + 1 行数据")
    header = [str(c).strip() for c in rows[0]]
    ncol = len(header)
    people = [r + [""] * (ncol - len(r)) for r in rows[1:] if any(str(c).strip() for c in r)]
    if not people:
        raise ValueError("没有找到员工数据行")

    wb = Workbook()
    ws = wb.active
    ws.title = "工资条"
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)

    r = 1
    block_rows = (3 if slip_title else 2) + 1  # 标题? + 表头 + 数据 + 空行
    for idx, person in enumerate(people):
        if slip_title:
            ws.cell(r, 1, slip_title).alignment = center
            ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=ncol)
            ws.cell(r, 1).font = Font(bold=True)
            r += 1
        for c in range(ncol):
            hc = ws.cell(r, c + 1, header[c])
            hc.alignment = center
            hc.border = _BORDER
            hc.font = Font(bold=True)
            hc.fill = _HEAD_FILL
            dc = ws.cell(r + 1, c + 1, person[c])
            dc.alignment = center
            dc.border = _BORDER
        r += 2
        # 每 per_page 条一个分页符
        if (idx + 1) % per_page == 0 and idx + 1 < len(people):
            ws.row_breaks.append(__import__("openpyxl").worksheet.pagebreak.Break(id=r - 1))
        r += 1  # 空行

    for c in range(ncol):
        ws.column_dimensions[get_column_letter(c + 1)].width = max(10, min(22, len(str(header[c])) * 2 + 6))

    ws.page_setup.orientation = "landscape" if ncol >= 7 else "portrait"
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    from openpyxl.worksheet.properties import PageSetupProperties

    ws.sheet_properties.pageSetUpPr = PageSetupProperties(fitToPage=True)

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# 多表合并 / 去重
# ---------------------------------------------------------------------------

def merge_sheets(files: list[bytes], dedupe: bool = False, key_column: str = "") -> bytes:
    """多个 Excel 合并成一个。默认认第一行是表头，用第一个文件的表头，后面文件的表头行跳过。
    dedupe=True 时去重：给了 key_column（表头名或列号）就按那一列去重，否则按整行去重。"""
    all_rows: list[list] = []
    header: list[str] | None = None

    for data in files:
        try:
            wb = load_workbook(io.BytesIO(data), data_only=True)
        except Exception:
            raise ValueError("有文件不是有效的 Excel（.xlsx）")
        for name in wb.sheetnames:
            ws = wb[name]
            sheet_rows = [[("" if c is None else c) for c in row] for row in ws.iter_rows(values_only=True)]
            sheet_rows = [r for r in sheet_rows if any(str(c).strip() for c in r)]
            if not sheet_rows:
                continue
            if header is None:
                header = [str(c).strip() for c in sheet_rows[0]]
                all_rows.extend(sheet_rows[1:])
            else:
                # 首行跟表头一样就当表头跳过，否则整张都是数据
                first = [str(c).strip() for c in sheet_rows[0]]
                all_rows.extend(sheet_rows[1:] if first == header else sheet_rows)

    if header is None:
        raise ValueError("没有读到任何数据")

    if dedupe:
        key_idx = None
        if key_column.strip():
            kc = key_column.strip()
            if kc.isdigit():
                key_idx = int(kc) - 1
            elif kc in header:
                key_idx = header.index(kc)
        seen = set()
        uniq = []
        for row in all_rows:
            k = (str(row[key_idx]).strip() if key_idx is not None and key_idx < len(row)
                 else tuple(str(c).strip() for c in row))
            if k in seen:
                continue
            seen.add(k)
            uniq.append(row)
        all_rows = uniq

    wb = Workbook()
    ws = wb.active
    ws.title = "合并结果"
    ws.append(header)
    for c in range(len(header)):
        ws.cell(1, c + 1).font = Font(bold=True)
    for row in all_rows:
        ws.append(list(row))

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
